"""
PPTX Export Service

Creates PowerPoint files with EDITABLE text boxes from Design JSON.
CRITICAL: Text must remain editable in PowerPoint, never baked into images.
"""

from typing import Dict, Any, Optional
import io
import base64
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.config import get_settings
from app.routers.jobs import update_job
from app.models.job import JobStatus


class PPTXExportService:
    """Service for exporting designs to PowerPoint with editable text."""

    # EMU conversion constants
    EMU_PER_INCH = 914400
    PX_PER_INCH = 96

    # Korean font mapping
    FONT_MAPPING = {
        "Pretendard": "맑은 고딕",
        "Noto Sans KR": "맑은 고딕",
        "Spoqa Han Sans": "맑은 고딕",
        "Inter": "Arial",
        "Roboto": "Arial",
    }

    @staticmethod
    def px_to_emu(px: float) -> int:
        """Convert pixels to EMU (English Metric Units)."""
        return int(px * PPTXExportService.EMU_PER_INCH / PPTXExportService.PX_PER_INCH)

    @staticmethod
    def hex_to_rgb(hex_color: str) -> tuple:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip("#")
        if len(hex_color) == 3:
            hex_color = "".join(c * 2 for c in hex_color)
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    @staticmethod
    async def export(design_json: Dict[str, Any], job_id: str) -> bytes:
        """
        Export Design JSON to PPTX file.

        CRITICAL: All text layers are exported as editable text boxes,
        NOT as images. This allows users to edit text in PowerPoint.
        """
        try:
            update_job(job_id, status=JobStatus.RUNNING, progress=10)

            # Create presentation
            prs = Presentation()

            # Get canvas dimensions
            canvas = design_json.get("canvas", {})
            width_px = canvas.get("width", 1080)
            height_px = canvas.get("height", 1080)

            # Set slide size
            prs.slide_width = Emu(PPTXExportService.px_to_emu(width_px))
            prs.slide_height = Emu(PPTXExportService.px_to_emu(height_px))

            update_job(job_id, progress=20)

            # Add blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            # Set background
            bg_color = canvas.get("backgroundColor", "#ffffff")
            PPTXExportService._set_slide_background(slide, bg_color)

            update_job(job_id, progress=30)

            # Process layers
            layers = design_json.get("layers", [])
            total_layers = len(layers)

            for i, layer in enumerate(layers):
                layer_type = layer.get("type")

                if layer_type == "background":
                    PPTXExportService._add_background_layer(slide, layer, width_px, height_px)
                elif layer_type == "text":
                    # CRITICAL: Add as editable text box
                    PPTXExportService._add_text_layer(slide, layer)
                elif layer_type == "shape":
                    PPTXExportService._add_shape_layer(slide, layer)
                elif layer_type in ["image", "product"]:
                    PPTXExportService._add_image_layer(slide, layer)

                progress = 30 + int((i + 1) / total_layers * 60)
                update_job(job_id, progress=progress)

            update_job(job_id, progress=95)

            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)

            update_job(
                job_id,
                status=JobStatus.COMPLETED,
                progress=100,
                output_data={"format": "pptx"},
            )

            return output.getvalue()

        except Exception as e:
            update_job(job_id, error_message=str(e))
            raise

    @staticmethod
    def _set_slide_background(slide, color: str):
        """Set solid color background for slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        r, g, b = PPTXExportService.hex_to_rgb(color)
        fill.fore_color.rgb = RgbColor(r, g, b)

    @staticmethod
    def _add_background_layer(slide, layer: Dict, canvas_w: int, canvas_h: int):
        """Add background layer (image or color)."""
        bg_image = layer.get("backgroundImage")
        if bg_image:
            # Add as image
            PPTXExportService._add_image_layer(
                slide,
                {
                    **layer,
                    "imageUrl": bg_image,
                    "position": {"x": 0, "y": 0},
                    "size": {"width": canvas_w, "height": canvas_h},
                },
            )
        else:
            # Just use color (already set by slide background)
            pass

    @staticmethod
    def _add_text_layer(slide, layer: Dict):
        """
        Add text as an EDITABLE text box.

        CRITICAL: This creates a native PowerPoint text box that users
        can edit directly. Text is NEVER rendered as an image.
        """
        position = layer.get("position", {})
        size = layer.get("size", {})

        left = Emu(PPTXExportService.px_to_emu(position.get("x", 0)))
        top = Emu(PPTXExportService.px_to_emu(position.get("y", 0)))
        width = Emu(PPTXExportService.px_to_emu(size.get("width", 200)))
        height = Emu(PPTXExportService.px_to_emu(size.get("height", 50)))

        # Add text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        # Configure text frame
        tf.word_wrap = True

        # Set vertical alignment
        vertical_align = layer.get("verticalAlign", "top")
        if vertical_align == "middle":
            tf.anchor = MSO_ANCHOR.MIDDLE
        elif vertical_align == "bottom":
            tf.anchor = MSO_ANCHOR.BOTTOM
        else:
            tf.anchor = MSO_ANCHOR.TOP

        # Add paragraph
        p = tf.paragraphs[0]
        p.text = layer.get("content", "")

        # Set alignment
        text_align = layer.get("textAlign", "left")
        if text_align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif text_align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        # Set font properties
        font = p.font
        font.size = Pt(layer.get("fontSize", 16))

        # Map font family for Korean compatibility
        font_family = layer.get("fontFamily", "Pretendard")
        font.name = PPTXExportService.FONT_MAPPING.get(font_family, font_family)

        # Set font weight (bold if >= 600)
        font.bold = layer.get("fontWeight", 400) >= 600

        # Set font style
        font.italic = layer.get("fontStyle", "normal") == "italic"

        # Set color
        color = layer.get("color", "#000000")
        r, g, b = PPTXExportService.hex_to_rgb(color)
        font.color.rgb = RgbColor(r, g, b)

        # Set text decoration
        text_decoration = layer.get("textDecoration", "none")
        if text_decoration == "underline":
            font.underline = True
        elif text_decoration == "line-through":
            # PowerPoint doesn't support strikethrough via python-pptx directly
            pass

        # Set opacity
        opacity = layer.get("opacity", 1)
        if opacity < 1:
            # Apply transparency to fill
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RgbColor(255, 255, 255)
            # Note: python-pptx has limited transparency support

    @staticmethod
    def _add_shape_layer(slide, layer: Dict):
        """Add shape layer."""
        from pptx.enum.shapes import MSO_SHAPE

        position = layer.get("position", {})
        size = layer.get("size", {})
        shape_type = layer.get("shapeType", "rectangle")

        left = Emu(PPTXExportService.px_to_emu(position.get("x", 0)))
        top = Emu(PPTXExportService.px_to_emu(position.get("y", 0)))
        width = Emu(PPTXExportService.px_to_emu(size.get("width", 100)))
        height = Emu(PPTXExportService.px_to_emu(size.get("height", 100)))

        # Map shape types
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "ellipse": MSO_SHAPE.OVAL,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        }
        mso_shape = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(mso_shape, left, top, width, height)

        # Set fill
        fill_color = layer.get("fill")
        if fill_color:
            shape.fill.solid()
            r, g, b = PPTXExportService.hex_to_rgb(fill_color)
            shape.fill.fore_color.rgb = RgbColor(r, g, b)
        else:
            shape.fill.background()

        # Set stroke
        stroke_color = layer.get("stroke")
        stroke_width = layer.get("strokeWidth", 0)
        if stroke_color and stroke_width > 0:
            shape.line.color.rgb = RgbColor(*PPTXExportService.hex_to_rgb(stroke_color))
            shape.line.width = Pt(stroke_width)
        else:
            shape.line.fill.background()

    @staticmethod
    def _add_image_layer(slide, layer: Dict):
        """Add image or product layer."""
        image_url = layer.get("imageUrl") or layer.get("backgroundImage")
        if not image_url:
            return

        position = layer.get("position", {})
        size = layer.get("size", {})

        left = Emu(PPTXExportService.px_to_emu(position.get("x", 0)))
        top = Emu(PPTXExportService.px_to_emu(position.get("y", 0)))
        width = Emu(PPTXExportService.px_to_emu(size.get("width", 200)))
        height = Emu(PPTXExportService.px_to_emu(size.get("height", 200)))

        try:
            if image_url.startswith("data:"):
                # Base64 data URL
                header, data = image_url.split(",", 1)
                image_bytes = base64.b64decode(data)
                image_stream = io.BytesIO(image_bytes)
            else:
                # HTTP URL - download image
                import httpx
                import asyncio

                async def download():
                    async with httpx.AsyncClient(timeout=30) as client:
                        response = await client.get(image_url)
                        return response.content

                # Run in event loop
                try:
                    loop = asyncio.get_event_loop()
                    image_bytes = loop.run_until_complete(download())
                except RuntimeError:
                    # No event loop - create one
                    image_bytes = asyncio.run(download())

                image_stream = io.BytesIO(image_bytes)

            slide.shapes.add_picture(image_stream, left, top, width, height)

        except Exception as e:
            print(f"Failed to add image: {e}")
            # Add placeholder rectangle instead
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RgbColor(200, 200, 200)
